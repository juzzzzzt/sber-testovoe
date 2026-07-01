from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from config import OUTPUT_DIR


def create_pptx(data, client_name):
    """Создание презентации в формате .pptx"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ЦВЕТОВАЯ СХЕМА СБЕРА (зелёная)
    SBER_GREEN = RGBColor(33, 149, 0)  # Основной зелёный #219500
    SBER_GREEN_DARK = RGBColor(5, 102, 0)  # Тёмно-зелёный для заголовков
    LIGHT_GRAY = RGBColor(240, 240, 240)
    WHITE = RGBColor(255, 255, 255)
    DARK_GRAY = RGBColor(80, 80, 80)

    def add_styled_slide(title_text, content_func):
        """Добавление слайда со стилизованным заголовком"""
        slide_layout = prs.slide_layouts[5]  # Пустой лейаут
        slide = prs.slides.add_slide(slide_layout)

        # Заголовок слайда (плашка) - ЗЕЛЁНЫЙ
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = SBER_GREEN
        shape.line.fill.background()

        tf = shape.text_frame
        tf.text = title_text
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.margin_left = Inches(0.5)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Контент
        content_func(slide)

    # Слайд 1: Титульный
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Полностью пустой
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = SBER_GREEN

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data["title_slide"]["title"]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = data["title_slide"]["subtitle"]
    p2.font.size = Pt(24)
    p2.font.color.rgb = LIGHT_GRAY
    p2.alignment = PP_ALIGN.CENTER

    # Слайд 2: Профиль
    def draw_profile(slide):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = True  # Авто-размер текста

        for i, bullet in enumerate(data["profile_slide"]["bullets"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(18)  # Уменьшили с 20 до 18
            p.font.name = 'Calibri'
            p.space_after = Pt(12)
            p.level = 0

    add_styled_slide(data["profile_slide"]["title"], draw_profile)

    # Слайд 3: Продукты
    # Слайд 3: Продукты (адаптивный размер шрифта)
    def draw_products(slide):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        num_products = len(data["products_slide"]["products"])

        # Адаптивный размер шрифта в зависимости от количества продуктов
        if num_products <= 3:
            title_size = Pt(22)
            desc_size = Pt(17)
        elif num_products <= 4:
            title_size = Pt(20)
            desc_size = Pt(16)
        else:
            title_size = Pt(18)
            desc_size = Pt(14)

        for i, prod in enumerate(data["products_slide"]["products"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"✔ {prod['name']}"
            p.font.size = title_size
            p.font.bold = True
            p.font.color.rgb = SBER_GREEN_DARK
            p.font.name = 'Calibri'

            p2 = tf.add_paragraph()
            p2.text = f"   {prod['reason']}"
            p2.font.size = desc_size
            p2.font.color.rgb = DARK_GRAY
            p2.font.name = 'Calibri'
            p2.space_after = Pt(12)

    add_styled_slide(data["products_slide"]["title"], draw_products)


    # Слайд 4: Новости
    def draw_news(slide):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, news in enumerate(data["news_slide"]["news_items"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"{news['title']}"
            p.font.size = Pt(18)  # Уменьшили с 22 до 18
            p.font.bold = True
            p.font.color.rgb = SBER_GREEN_DARK
            p.font.name = 'Calibri'

            p2 = tf.add_paragraph()
            p2.text = f"   {news['summary']}"
            p2.font.size = Pt(15)  # Уменьшили с 18 до 15
            p2.font.color.rgb = DARK_GRAY
            p2.font.name = 'Calibri'
            p2.space_after = Pt(12)

    add_styled_slide(data["news_slide"]["title"], draw_news)

    # Сохранение файла
    filename = f"presentation_{client_name.strip().replace(' ', '_')}.pptx"
    filepath = f"{OUTPUT_DIR}/{filename}"
    prs.save(filepath)
    print(f"Презентация сохранена: {filepath}")

    return filepath